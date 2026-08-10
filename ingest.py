"""Document ingestion. [Track B]

Turn an uploaded PDF / Word / Markdown / PowerPoint into decision memory:
  1. extract_text  — pull raw text out of the file (per format)
  2. extract_key_points — Gemini reads it and returns the decisions it records
  3. ingest_file — append each of those to memory via Track A's record_decision

Deps (install yourself):  pip install pypdf python-docx python-pptx
"""
import io
import os
import re
import sys

from dotenv import load_dotenv

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "db"))
import tools as memory  # db/tools.py — record_decision

import distiller  # first-stage agent: raw text -> compressed per-change records

load_dotenv()


def extract_text(name: str, data: bytes) -> str:
    """Extract plain text from a file's bytes, dispatched by extension."""
    ext = name.lower().rsplit(".", 1)[-1] if "." in name else ""

    if ext == "pdf":
        try:
            from pypdf import PdfReader
        except ImportError:
            raise RuntimeError("PDF support needs pypdf — run: pip install pypdf")
        reader = PdfReader(io.BytesIO(data))
        return "\n".join((page.extract_text() or "") for page in reader.pages)

    if ext == "docx":
        try:
            import docx
        except ImportError:
            raise RuntimeError("Word support needs python-docx — run: pip install python-docx")
        document = docx.Document(io.BytesIO(data))
        return "\n".join(p.text for p in document.paragraphs)

    if ext == "pptx":
        try:
            from pptx import Presentation
        except ImportError:
            raise RuntimeError("PowerPoint support needs python-pptx — run: pip install python-pptx")
        prs = Presentation(io.BytesIO(data))
        chunks = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    chunks.append(shape.text_frame.text)
        return "\n".join(chunks)

    if ext in ("md", "markdown", "txt"):
        return data.decode("utf-8", errors="ignore")

    raise ValueError(f"Unsupported file type: .{ext} (use pdf, docx, md, txt, or pptx)")


_DATE_RE = re.compile(r"^\d{4}-\d{2}-\d{2}$")


def _parse_date(value):
    """Accept only a clean YYYY-MM-DD string; anything else -> None."""
    if isinstance(value, str) and _DATE_RE.match(value.strip()):
        return value.strip()
    return None


def _known_topics() -> list[str]:
    """Distinct topics already in memory, so the distiller can reuse exact labels."""
    try:
        return sorted({d["topic"] for d in memory.list_decisions()})
    except Exception:
        return []


def ingest_file(name: str, data: bytes) -> list[dict]:
    """Full pipeline for one file: extract text -> distiller agent -> append each change.
    Returns the items actually recorded (topic + new_state)."""
    text = extract_text(name, data)
    changes = distiller.distill(text, _known_topics())
    recorded = []
    for c in changes:
        topic = (c.get("topic") or "").strip()
        new_state = (c.get("new_state") or "").strip()
        if not topic or not new_state:
            continue
        memory.record_decision(
            topic=topic,
            old_state=(c.get("old_state") or None),
            new_state=new_state,
            cause=c.get("cause"),
            trigger_event=c.get("trigger_event"),
            tension=c.get("tension"),
            recorded_by=f"Ingest: {name}",
            created_at=_parse_date(c.get("date")),
        )
        recorded.append({"topic": topic, "new_state": new_state})
    return recorded
